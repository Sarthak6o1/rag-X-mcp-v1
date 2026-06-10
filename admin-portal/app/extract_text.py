"""
Extract plain text to match what ``POST /api/ingest`` expects.
Mirror behavior of a bulk run on ``backend`` (same page_index/trees, same extraction choices).

**Keep in sync** with your backend pipeline when you change supported types or extractors.
"""
from __future__ import annotations

import csv
import json
from io import StringIO, BytesIO
from pathlib import Path

from app import config
from app.page_index import MEDIA_EXTENSIONS

SUPPORTED = (
    {".txt", ".md", ".pdf", ".docx", ".xlsx", ".pptx", ".csv", ".json"}
    | MEDIA_EXTENSIONS
)

SUPPORTED_LABEL = (
    ".txt, .md, .pdf, .docx, .xlsx, .pptx, .csv, .json, "
    + ", ".join(sorted(MEDIA_EXTENSIONS))
)

# HTML file input + backend parity: same set as ``page_index.MEDIA_EXTENSIONS`` and document types in backend bulk flow.
UPLOAD_ACCEPT = ",".join(sorted(SUPPORTED))


def extract_text(
    data: bytes,
    filename: str,
    work_path: Path | None = None,
) -> tuple[str, str | None]:
    """
    Returns (text, error). ``work_path`` is the temp file for media sidecar
    resolution (``file.ext.txt`` / ``_transcripts/…``), as in ``page_index``.
    """
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED:
        return "", f"Unsupported type {suffix!r}. Allowed: {', '.join(sorted(SUPPORTED))}."
    if not data and suffix not in MEDIA_EXTENSIONS:
        return "", "Empty file."
    if suffix in MEDIA_EXTENSIONS:
        if work_path is None:
            return "", "Internal error: media path missing."
        # Optional sidecar (same as bulk / backend); if missing, use speech-to-text when enabled.
        text, _ = _media_via_transcript(work_path)
        if text.strip():
            return text, None
        if config.TRANSCRIBE_MEDIA:
            from app.transcribe_media import transcribe_file

            t2, err2 = transcribe_file(work_path)
            if t2.strip():
                return t2, None
            return "", err2 or "We couldn’t get usable text from this video. It may have no clear speech, or the file may be too low quality. Try again or use a different file."
        return "", (
            "This app isn’t set up to turn video into text automatically. Ask your team to enable that, or add a plain text transcript with the same name as the file."
        )
    if not data or not data.strip():
        return "", "Empty file."
    if suffix in {".txt", ".md"}:
        return _plain_text(data)
    if suffix == ".json":
        return _json_text(data)
    if suffix == ".csv":
        return _csv_text(data)
    if suffix == ".pdf":
        return _pdf_text(data)
    if suffix == ".docx":
        return _docx_text(data)
    if suffix == ".xlsx":
        return _xlsx_text(data)
    if suffix == ".pptx":
        return _pptx_text(data)
    return "", "Internal error: unhandled type."


def _media_via_transcript(path: Path) -> tuple[str, str | None]:
    candidates = [
        Path(f"{path}.txt"),
        path.parent / "_transcripts" / f"{path.name}.txt",
    ]
    for candidate in candidates:
        if not candidate.is_file():
            continue
        try:
            text = candidate.read_text(encoding="utf-8", errors="replace")
        except OSError as exc:
            return "", f"Could not read transcript: {exc}"
        if text.strip():
            return text, None
    return "", None


def _plain_text(data: bytes) -> tuple[str, str | None]:
    for enc in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            text = data.decode(enc)
            if text.strip():
                return text, None
        except UnicodeDecodeError:
            continue
    return "", "Could not decode text file."


def _json_text(data: bytes) -> tuple[str, str | None]:
    raw, err = _plain_text(data)
    if err:
        return raw, err
    try:
        obj = json.loads(raw)
    except json.JSONDecodeError as exc:
        return "", f"Invalid JSON: {exc}"
    try:
        text = json.dumps(obj, ensure_ascii=False, indent=2)
    except (TypeError, ValueError) as exc:
        return "", f"JSON re-serialization failed: {exc}"
    return text, None


def _csv_text(data: bytes) -> tuple[str, str | None]:
    raw, err = _plain_text(data)
    if err:
        return raw, err
    out_lines: list[str] = []
    f = StringIO(raw)
    try:
        reader = csv.reader(f)
        for row in reader:
            out_lines.append("\t".join(str(c) for c in row if c is not None))
    except csv.Error as exc:
        return "", f"CSV parse error: {exc}"
    text = "\n".join(out_lines).strip()
    if not text:
        return "", "CSV produced no cell content."
    return text, None


def _pdf_text(data: bytes) -> tuple[str, str | None]:
    try:
        from pypdf import PdfReader
    except ImportError:
        return "", "pypdf not available."
    try:
        reader = PdfReader(BytesIO(data))
        parts: list[str] = []
        for page in reader.pages:
            t = page.extract_text()
            if t:
                parts.append(t)
        text = "\n\n".join(parts).strip()
        if not text:
            return "", "No extractable text in PDF (may be scanned image-only)."
        return text, None
    except Exception as exc:  # noqa: BLE001
        return "", f"PDF read error: {exc}"


def _docx_text(data: bytes) -> tuple[str, str | None]:
    try:
        import docx
    except ImportError:
        return "", "python-docx not available."
    try:
        document = docx.Document(BytesIO(data))
        parts = [p.text for p in document.paragraphs if p.text and p.text.strip()]
        text = "\n\n".join(parts).strip()
        if not text:
            return "", "No text in DOCX."
        return text, None
    except Exception as exc:  # noqa: BLE001
        return "", f"DOCX read error: {exc}"


def _xlsx_text(data: bytes) -> tuple[str, str | None]:
    try:
        import openpyxl
    except ImportError:
        return "", "openpyxl not available."
    try:
        wb = openpyxl.load_workbook(BytesIO(data), read_only=True, data_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"# {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(cell) for cell in row if cell is not None and str(cell).strip()]
                if cells:
                    parts.append("\t".join(cells))
        wb.close()
        text = "\n".join(parts).strip()
        if not text:
            return "", "Workbook is empty."
        return text, None
    except Exception as exc:  # noqa: BLE001
        return "", f"XLSX read error: {exc}"


def _pptx_text(data: bytes) -> tuple[str, str | None]:
    try:
        from pptx import Presentation
    except ImportError:
        return "", "python-pptx not available."
    try:
        prs = Presentation(BytesIO(data))
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs)
                        if text.strip():
                            parts.append(text.strip())
        out = "\n\n".join(parts).strip()
        if not out:
            return "", "No text in PPTX."
        return out, None
    except Exception as exc:  # noqa: BLE001
        return "", f"PPTX read error: {exc}"


# Alias for clarity if you name it like the backend
extract_for_ingest = extract_text
