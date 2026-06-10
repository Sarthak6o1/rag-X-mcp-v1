"""Mirror of ``backend/rag_pipeline/page_index.py`` — keep in sync manually (admin-only copy)."""
from __future__ import annotations

import re
from pathlib import Path
from typing import Any

MEDIA_EXTENSIONS = {".mp4", ".mov", ".mkv", ".avi", ".wmv", ".webm"}


def _node(node_id: str, parent_id: str | None, title: str, level: int) -> dict[str, Any]:
    return {"id": node_id, "parent_id": parent_id, "title": title, "level": level}


def _tree_from_titles(titles: list[tuple[int, str]], root_title: str) -> dict[str, Any]:
    nodes: list[dict[str, Any]] = [_node("n0", None, root_title, 0)]
    stack: list[tuple[str, int]] = [("n0", 0)]
    next_id = 1
    for level, title in titles:
        if not title.strip():
            continue
        while stack and stack[-1][1] >= level:
            stack.pop()
        parent_id = stack[-1][0] if stack else "n0"
        nid = f"n{next_id}"
        next_id += 1
        nodes.append(_node(nid, parent_id, title.strip()[:500], level))
        stack.append((nid, level))
    return {"root_id": "n0", "nodes": nodes}


def _headings_docx(path: Path) -> list[tuple[int, str]] | None:
    try:
        import docx
    except ImportError:
        return None
    try:
        document = docx.Document(str(path))
    except (OSError, ValueError):
        return None
    out: list[tuple[int, str]] = []
    for p in document.paragraphs:
        text = (p.text or "").strip()
        if not text:
            continue
        style = (p.style.name if p.style else "") or ""
        if style.startswith("Heading"):
            try:
                lvl = int(style.replace("Heading", "").strip() or "1")
            except ValueError:
                lvl = 1
            out.append((min(lvl, 6), text))
    return out or None


def _outline_pdf(path: Path) -> list[tuple[int, str]] | None:
    try:
        from pypdf import PdfReader
    except ImportError:
        return None
    try:
        reader = PdfReader(str(path))
        outline = reader.outline
    except (OSError, ValueError):
        return None

    titles: list[tuple[int, str]] = []

    def walk(items: object, depth: int) -> None:
        if items is None:
            return
        if isinstance(items, list):
            for item in items:
                walk(item, depth)
            return
        title = getattr(items, "title", None)
        if title:
            titles.append((min(depth + 1, 6), str(title)))
        children = getattr(items, "children", None)
        if children:
            walk(children, depth + 1)

    if not outline:
        return None
    try:
        walk(outline, 0)
    except (TypeError, AttributeError):
        return None
    return titles or None


def _headings_markdown(path: Path) -> list[tuple[int, str]] | None:
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
    except OSError:
        return None
    out: list[tuple[int, str]] = []
    for line in text.splitlines():
        stripped = line.strip()
        if stripped.startswith("#"):
            level = 0
            for ch in stripped:
                if ch == "#":
                    level += 1
                else:
                    break
            title = stripped[level:].strip()
            if title:
                out.append((min(level, 6), title))
    return out or None


def _sheets_xlsx(path: Path) -> list[tuple[int, str]] | None:
    try:
        import openpyxl
    except ImportError:
        return None
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        names = wb.sheetnames
        wb.close()
    except (OSError, ValueError):
        return None
    if not names:
        return None
    return [(1, name) for name in names if name]


def _pdf_page_count(path: Path) -> int | None:
    try:
        from pypdf import PdfReader
    except ImportError:
        return None
    try:
        reader = PdfReader(str(path))
        return len(reader.pages)
    except (OSError, ValueError):
        return None


def _slides_pptx(path: Path) -> list[tuple[int, str]] | None:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
    except Exception:
        return None
    out: list[tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()
        out.append((1, title or f"Slide {i}"))
    return out or None


def _sections_from_transcript(path: Path) -> list[tuple[int, str]] | None:
    """Extract headings from a sidecar transcript file (timestamps become sections)."""
    candidates = [
        Path(f"{path}.txt"),
        path.parent / "_transcripts" / f"{path.name}.txt",
    ]
    for candidate in candidates:
        if not candidate.is_file():
            continue
        try:
            text = candidate.read_text(encoding="utf-8", errors="replace")
        except OSError:
            continue
        headings = _headings_markdown(candidate)
        if headings:
            return headings
        lines = text.splitlines()
        ts_sections: list[tuple[int, str]] = []
        for line in lines:
            stripped = line.strip()
            if re.match(r"^\[?\d{1,2}:\d{2}", stripped):
                ts_sections.append((1, stripped[:120]))
        if ts_sections:
            return ts_sections
    return None


def build_document_tree(path: Path) -> dict[str, Any]:
    """
    Local PageIndex-style hierarchical index: TOC / headings / sheets / slides + metadata.
    """
    suffix = path.suffix.lower()
    root_title = path.stem[:200] or path.name
    meta: dict[str, Any] = {"schema_version": "local-pageindex-v2", "source": "rag_pipeline.page_index"}

    titles: list[tuple[int, str]] | None = None
    if suffix == ".docx":
        titles = _headings_docx(path)
    elif suffix == ".pdf":
        titles = _outline_pdf(path)
        pc = _pdf_page_count(path)
        if pc is not None:
            meta["pdf_pages"] = pc
    elif suffix in {".md", ".txt"}:
        titles = _headings_markdown(path)
    elif suffix == ".xlsx":
        titles = _sheets_xlsx(path)
    elif suffix == ".pptx":
        titles = _slides_pptx(path)
    elif suffix in MEDIA_EXTENSIONS:
        titles = _sections_from_transcript(path)

    if titles:
        tree = _tree_from_titles(titles, root_title)
        tree["meta"] = meta
        return tree

    tree = {
        "root_id": "n0",
        "nodes": [_node("n0", None, root_title, 0)],
        "meta": meta,
    }
    if suffix == ".pdf":
        pc = _pdf_page_count(path)
        if pc is not None:
            tree["meta"]["pdf_pages"] = pc
    return tree
